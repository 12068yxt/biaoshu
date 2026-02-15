#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT导入工具
将生成的SVG文件转换为PPT幻灯片
"""

import os
import logging
from typing import List
from datetime import datetime

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTImporter:
    def __init__(self):
        """初始化PPT导入器"""
        self.svg_dir = "generated_svgs"
        self.ppt_dir = "ppt_presentations"
        os.makedirs(self.ppt_dir, exist_ok=True)
    
    def get_svg_files(self) -> List[str]:
        """获取所有SVG文件"""
        if not os.path.exists(self.svg_dir):
            logger.error(f"SVG目录不存在: {self.svg_dir}")
            return []
        
        svg_files = [f for f in os.listdir(self.svg_dir) if f.endswith('.svg')]
        svg_files.sort()
        return svg_files
    
    def convert_svg_to_image(self, svg_path: str) -> str:
        """将SVG转换为图片格式（PNG）"""
        try:
            import cairosvg
            
            # 生成PNG文件路径
            png_path = svg_path.replace('.svg', '.png')
            
            # 转换SVG到PNG
            cairosvg.svg2png(
                url=svg_path,
                write_to=png_path,
                output_width=1920,
                output_height=1080
            )
            
            logger.info(f"转换完成: {os.path.basename(png_path)}")
            return png_path
            
        except ImportError:
            logger.warning("未安装cairosvg，跳过SVG转换")
            return ""
        except Exception as e:
            logger.error(f"转换失败 {svg_path}: {str(e)}")
            return ""
    
    def create_powerpoint_presentation(self, image_files: List[str]) -> str:
        """创建PowerPoint演示文稿"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            
            # 创建新的演示文稿
            prs = Presentation()
            
            # 设置幻灯片尺寸为16:9宽屏
            prs.slide_width = Inches(13.33)  # 1920px
            prs.slide_height = Inches(7.5)   # 1080px
            
            for img_path in image_files:
                if not os.path.exists(img_path):
                    continue
                
                # 添加新幻灯片（空白布局）
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 添加图片（填满整个幻灯片）
                left = top = Inches(0)
                width = prs.slide_width
                height = prs.slide_height
                
                slide.shapes.add_picture(img_path, left, top, width, height)
                
                # 添加页脚信息
                footer_textbox = slide.shapes.add_textbox(
                    Inches(0.5), prs.slide_height - Inches(0.8), 
                    prs.slide_width - Inches(1), Inches(0.7)
                )
                footer_frame = footer_textbox.text_frame
                footer_frame.text = f"技术文档可视化 - {datetime.now().strftime('%Y-%m-%d')}"
                footer_para = footer_frame.paragraphs[0]
                footer_para.font.size = Pt(14)
                footer_para.font.color.rgb = RGBColor(100, 100, 100)
                
                logger.info(f"添加幻灯片: {os.path.basename(img_path)}")
            
            # 保存演示文稿
            ppt_filename = f"technical_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            ppt_path = os.path.join(self.ppt_dir, ppt_filename)
            prs.save(ppt_path)
            
            logger.info(f"PPT文件已创建: {ppt_path}")
            return ppt_path
            
        except ImportError:
            logger.error("请安装python-pptx库: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"创建PPT失败: {str(e)}")
            return ""
    
    def create_html_presentation(self, svg_files: List[str]) -> str:
        """创建HTML格式的演示文稿（备选方案）"""
        try:
            html_content = '''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>技术文档演示</title>
    <style>
        body {
            margin: 0;
            padding: 0;
            font-family: Arial, sans-serif;
            background: #f0f0f0;
        }
        .slide-container {
            width: 100vw;
            height: 100vh;
            display: flex;
            justify-content: center;
            align-items: center;
            position: relative;
        }
        .slide {
            width: 1920px;
            height: 1080px;
            box-shadow: 0 0 20px rgba(0,0,0,0.3);
            display: none;
        }
        .slide.active {
            display: block;
        }
        .controls {
            position: fixed;
            bottom: 20px;
            left: 50%;
            transform: translateX(-50%);
            z-index: 1000;
        }
        .btn {
            background: #007bff;
            color: white;
            border: none;
            padding: 10px 20px;
            margin: 0 10px;
            border-radius: 5px;
            cursor: pointer;
            font-size: 16px;
        }
        .btn:hover {
            background: #0056b3;
        }
        .slide-number {
            position: fixed;
            top: 20px;
            right: 20px;
            background: rgba(0,0,0,0.7);
            color: white;
            padding: 5px 15px;
            border-radius: 15px;
            font-size: 14px;
        }
    </style>
</head>
<body>
    <div class="slide-number" id="slideNumber">1 / {total_slides}</div>
    
    <div class="controls">
        <button class="btn" onclick="prevSlide()">上一页</button>
        <button class="btn" onclick="nextSlide()">下一页</button>
    </div>
'''

            # 添加每张幻灯片
            for i, svg_file in enumerate(svg_files, 1):
                svg_path = os.path.join(self.svg_dir, svg_file)
                display_style = "active" if i == 1 else ""
                html_content += f'''
    <div class="slide-container">
        <div class="slide {display_style}" id="slide{i}">
            <img src="{svg_path}" alt="Slide {i}" style="width: 100%; height: 100%; object-fit: cover;">
        </div>
    </div>'''

            html_content += '''
    <script>
        let currentSlide = 1;
        const totalSlides = {total_slides};
        
        function showSlide(n) {
            // 隐藏所有幻灯片
            document.querySelectorAll('.slide').forEach(slide => {
                slide.classList.remove('active');
            });
            
            // 显示当前幻灯片
            if (n > totalSlides) currentSlide = 1;
            if (n < 1) currentSlide = totalSlides;
            
            document.getElementById(`slide${currentSlide}`).classList.add('active');
            document.getElementById('slideNumber').textContent = `${currentSlide} / ${totalSlides}`;
        }
        
        function nextSlide() {
            currentSlide++;
            showSlide(currentSlide);
        }
        
        function prevSlide() {
            currentSlide--;
            showSlide(currentSlide);
        }
        
        // 键盘控制
        document.addEventListener('keydown', function(event) {
            if (event.key === 'ArrowRight' || event.key === ' ') {
                nextSlide();
            } else if (event.key === 'ArrowLeft') {
                prevSlide();
            }
        });
    </script>
</body>
</html>'''

            html_content = html_content.replace('{total_slides}', str(len(svg_files)))
            
            # 保存HTML文件
            html_filename = f"technical_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
            html_path = os.path.join(self.ppt_dir, html_filename)
            
            with open(html_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            logger.info(f"HTML演示文稿已创建: {html_path}")
            return html_path
            
        except Exception as e:
            logger.error(f"创建HTML演示文稿失败: {str(e)}")
            return ""
    
    def import_to_ppt(self):
        """执行导入流程"""
        # 获取SVG文件
        svg_files = self.get_svg_files()
        if not svg_files:
            logger.error("未找到SVG文件")
            return
        
        logger.info(f"找到 {len(svg_files)} 个SVG文件")
        
        # 方法1: 尝试创建PowerPoint文件
        image_files = []
        for svg_file in svg_files:
            svg_path = os.path.join(self.svg_dir, svg_file)
            png_path = self.convert_svg_to_image(svg_path)
            if png_path:
                image_files.append(png_path)
        
        ppt_path = ""
        if image_files:
            ppt_path = self.create_powerpoint_presentation(image_files)
        
        # 方法2: 创建HTML演示文稿（总是可用）
        html_path = self.create_html_presentation(svg_files)
        
        # 输出结果
        logger.info("\n=== 导入完成 ===")
        if ppt_path:
            logger.info(f"PowerPoint文件: {ppt_path}")
        logger.info(f"HTML演示文稿: {html_path}")
        logger.info(f"输出目录: {os.path.abspath(self.ppt_dir)}")

def main():
    """主函数"""
    importer = PPTImporter()
    importer.import_to_ppt()

if __name__ == "__main__":
    main()