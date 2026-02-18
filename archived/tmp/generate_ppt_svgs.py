#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT风格SVG生成器
为技术文档章节生成兼容PPT的SVG图表，并自动导入到PPT中
"""

import os
import json
import re
import asyncio
import aiohttp
from typing import List, Dict, Tuple
import xml.etree.ElementTree as ET
from datetime import datetime
import logging

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('ppt_svg_generation.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class PPTSVGGenerator:
    def __init__(self):
        """初始化SVG生成器"""
        self.api_key = os.getenv('ALIYUN_API_KEY')
        if not self.api_key:
            raise ValueError("请设置环境变量 ALIYUN_API_KEY")
        
        self.api_url = "https://dashscope.aliyuncs.com/api/v1/services/aigc/text-generation/generation"
        self.headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }
        
        # 创建输出目录
        self.svg_output_dir = "ppt_svgs"
        self.ppt_output_dir = "ppt_presentations"
        os.makedirs(self.svg_output_dir, exist_ok=True)
        os.makedirs(self.ppt_output_dir, exist_ok=True)
        
        # 进度文件
        self.progress_file = "svg_generation_progress.json"
        
    def load_progress(self) -> Dict:
        """加载处理进度"""
        if os.path.exists(self.progress_file):
            with open(self.progress_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {"completed": [], "failed": []}
    
    def save_progress(self, progress: Dict):
        """保存处理进度"""
        with open(self.progress_file, 'w', encoding='utf-8') as f:
            json.dump(progress, f, ensure_ascii=False, indent=2)
    
    def get_txt_files(self) -> List[str]:
        """获取所有txt文件"""
        txt_files = [f for f in os.listdir('.') if f.endswith('.txt') and f.startswith('chapter_')]
        txt_files.sort()
        logger.info(f"找到 {len(txt_files)} 个章节文件")
        return txt_files
    
    def extract_chapter_info(self, filename: str) -> Tuple[str, str]:
        """从文件名提取章节编号和标题"""
        # 匹配格式: chapter_01_1.2.1.3.1.1 标题.txt
        pattern = r'chapter_(\d+)_(.+?)\s+(.+?)\.txt'
        match = re.match(pattern, filename)
        if match:
            chapter_num = match.group(1)
            section_id = match.group(2)
            title = match.group(3)
            return chapter_num, f"{section_id} {title}"
        else:
            # 备用匹配
            pattern2 = r'chapter_(\d+)_(.+?)\.txt'
            match2 = re.match(pattern2, filename)
            if match2:
                chapter_num = match2.group(1)
                title = match2.group(2)
                return chapter_num, title
        return "0", filename.replace('.txt', '')
    
    async def generate_svg_content(self, chapter_title: str, content_summary: str) -> str:
        """使用Qwen coder模型生成SVG代码"""
        prompt = f"""
请为以下技术章节内容生成一个专业的PPT风格SVG图表代码：

章节标题: {chapter_title}
内容摘要: {content_summary[:500]}...

要求：
1. 生成标准的SVG代码，可以直接在PPT中使用
2. 图表风格要专业、简洁，适合技术演示
3. 包含标题区域和主要内容展示区域
4. 使用现代化的配色方案（蓝白灰为主）
5. 添加适当的图标和装饰元素
6. 尺寸为1920x1080像素（PPT宽屏比例）
7. 确保文字清晰可读
8. 返回纯SVG代码，不要包含其他说明文字

请直接返回SVG代码，不需要额外解释。
"""

        payload = {
            "model": "qwen-plus",
            "input": {
                "messages": [
                    {
                        "role": "user",
                        "content": prompt
                    }
                ]
            },
            "parameters": {
                "temperature": 0.7,
                "max_tokens": 2000
            }
        }

        async with aiohttp.ClientSession() as session:
            try:
                async with session.post(
                    self.api_url,
                    headers=self.headers,
                    json=payload,
                    timeout=aiohttp.ClientTimeout(total=120)
                ) as response:
                    if response.status == 200:
                        result = await response.json()
                        svg_code = result['output']['text'].strip()
                        
                        # 清理SVG代码，确保格式正确
                        svg_code = self.clean_svg_code(svg_code)
                        return svg_code
                    else:
                        error_text = await response.text()
                        logger.error(f"API调用失败: {response.status} - {error_text}")
                        return self.generate_default_svg(chapter_title, content_summary)
            except Exception as e:
                logger.error(f"生成SVG时出错: {str(e)}")
                return self.generate_default_svg(chapter_title, content_summary)
    
    def clean_svg_code(self, svg_code: str) -> str:
        """清理SVG代码，确保格式正确"""
        # 移除可能的markdown代码块标记
        svg_code = svg_code.replace('```xml', '').replace('```svg', '').replace('```', '')
        svg_code = svg_code.strip()
        
        # 如果没有SVG标签，添加默认包装
        if not svg_code.startswith('<svg'):
            svg_code = f'<svg xmlns="http://www.w3.org/2000/svg" width="1920" height="1080" viewBox="0 0 1920 1080">{svg_code}</svg>'
        
        return svg_code
    
    def generate_default_svg(self, title: str, content: str) -> str:
        """生成默认的SVG模板"""
        # 截取标题和内容用于显示
        short_title = title[:50] + "..." if len(title) > 50 else title
        short_content = content[:200] + "..." if len(content) > 200 else content
        
        svg_template = f'''<svg xmlns="http://www.w3.org/2000/svg" width="1920" height="1080" viewBox="0 0 1920 1080">
  <!-- 背景 -->
  <rect width="1920" height="1080" fill="#f8f9fa"/>
  
  <!-- 标题区域背景 -->
  <rect x="100" y="100" width="1720" height="150" rx="15" fill="#1e3a8a" opacity="0.9"/>
  
  <!-- 标题文字 -->
  <text x="150" y="190" font-family="Arial, sans-serif" font-size="48" fill="white" font-weight="bold">
    {short_title}
  </text>
  
  <!-- 内容区域背景 -->
  <rect x="100" y="280" width="1720" height="600" rx="15" fill="white" stroke="#e2e8f0" stroke-width="2"/>
  
  <!-- 内容文字 -->
  <text x="150" y="330" font-family="Arial, sans-serif" font-size="24" fill="#334155">
    {short_content}
  </text>
  
  <!-- 装饰元素 -->
  <circle cx="1800" cy="175" r="20" fill="#3b82f6"/>
  <circle cx="1850" cy="175" r="15" fill="#60a5fa"/>
  <circle cx="1800" cy="900" r="25" fill="#1e3a8a" opacity="0.3"/>
  <circle cx="150" cy="950" r="15" fill="#3b82f6" opacity="0.5"/>
</svg>'''
        
        return svg_template
    
    def save_svg_file(self, chapter_num: str, title: str, svg_content: str) -> str:
        """保存SVG文件"""
        # 创建安全的文件名
        safe_title = re.sub(r'[<>:"/\\|?*]', '_', title)[:50]
        filename = f"chapter_{chapter_num}_{safe_title}.svg"
        filepath = os.path.join(self.svg_output_dir, filename)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(svg_content)
        
        logger.info(f"已保存SVG文件: {filename}")
        return filepath
    
    def create_ppt_from_svgs(self, svg_files: List[Dict]):
        """创建PPT文件（这里提供思路，实际需要python-pptx库）"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation()
            prs.slide_width = Inches(13.33)  # 1920px
            prs.slide_height = Inches(7.5)   # 1080px
            
            for svg_info in svg_files:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
                
                # 添加标题
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
                title_frame = title_box.text_frame
                title_frame.text = svg_info['title']
                title_frame.paragraphs[0].font.size = 360000  # 36pt
                title_frame.paragraphs[0].font.bold = True
                
                # 这里需要将SVG转换为图片格式才能插入PPT
                # 实际实现需要额外的转换步骤
                
            ppt_filename = f"technical_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            ppt_filepath = os.path.join(self.ppt_output_dir, ppt_filename)
            prs.save(ppt_filepath)
            
            logger.info(f"PPT文件已创建: {ppt_filepath}")
            return ppt_filepath
            
        except ImportError:
            logger.warning("未安装python-pptx库，无法直接生成PPT文件")
            logger.info("请安装: pip install python-pptx")
            return None
        except Exception as e:
            logger.error(f"创建PPT时出错: {str(e)}")
            return None
    
    async def process_single_chapter(self, filename: str, semaphore: asyncio.Semaphore) -> Dict:
        """处理单个章节文件"""
        async with semaphore:
            try:
                chapter_num, title = self.extract_chapter_info(filename)
                
                # 读取文件内容
                with open(filename, 'r', encoding='utf-8') as f:
                    content = f.read()
                
                # 生成内容摘要（取前1000字符）
                content_summary = content[:1000] if len(content) > 1000 else content
                
                logger.info(f"正在处理章节 {chapter_num}: {title}")
                
                # 生成SVG
                svg_content = await self.generate_svg_content(title, content_summary)
                
                # 保存SVG文件
                svg_filepath = self.save_svg_file(chapter_num, title, svg_content)
                
                return {
                    'filename': filename,
                    'chapter_num': chapter_num,
                    'title': title,
                    'svg_path': svg_filepath,
                    'status': 'success'
                }
                
            except Exception as e:
                logger.error(f"处理文件 {filename} 时出错: {str(e)}")
                return {
                    'filename': filename,
                    'status': 'failed',
                    'error': str(e)
                }
    
    async def generate_all_svgs(self, max_concurrent: int = 3):
        """批量生成所有章节的SVG"""
        txt_files = self.get_txt_files()
        progress = self.load_progress()
        
        # 过滤已处理的文件
        remaining_files = [f for f in txt_files if f not in progress['completed']]
        
        if not remaining_files:
            logger.info("所有文件已处理完成")
            return
        
        logger.info(f"开始处理 {len(remaining_files)} 个剩余文件")
        
        semaphore = asyncio.Semaphore(max_concurrent)
        tasks = [self.process_single_chapter(filename, semaphore) for filename in remaining_files]
        
        results = []
        for i, task in enumerate(asyncio.as_completed(tasks)):
            result = await task
            results.append(result)
            
            # 更新进度
            if result['status'] == 'success':
                progress['completed'].append(result['filename'])
            else:
                progress['failed'].append(result['filename'])
            
            self.save_progress(progress)
            
            logger.info(f"进度: {i+1}/{len(remaining_files)} 完成")
        
        # 创建PPT（如果可能）
        successful_results = [r for r in results if r['status'] == 'success']
        if successful_results:
            self.create_ppt_from_svgs(successful_results)
        
        # 输出统计
        success_count = len([r for r in results if r['status'] == 'success'])
        failed_count = len([r for r in results if r['status'] == 'failed'])
        
        logger.info(f"处理完成! 成功: {success_count}, 失败: {failed_count}")

def main():
    """主函数"""
    try:
        generator = PPTSVGGenerator()
        asyncio.run(generator.generate_all_svgs(max_concurrent=3))
    except KeyboardInterrupt:
        logger.info("用户中断处理")
    except Exception as e:
        logger.error(f"程序执行出错: {str(e)}")

if __name__ == "__main__":
    main()